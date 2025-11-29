import os
import pypdf
from docx import Document
from pptx import Presentation

def extract_text_from_file(file_path):
    """
    Extracts text from PDF, DOCX, PPTX, or TXT files.
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext == '.txt':
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def extract_text_from_pdf(file_path):
    text = ""
    try:
        with open(file_path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error reading PDF: {e}")
    return text

def extract_text_from_docx(file_path):
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error reading DOCX: {e}")
    return text

def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX: {e}")
    return text

def extract_text_from_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"Error reading TXT: {e}")
        return ""
